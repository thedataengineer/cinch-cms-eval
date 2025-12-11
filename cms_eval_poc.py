"""
Home Warranty CMS Evaluation Framework – Streamlit PoC
Interactive tool to evaluate and score CMS platforms for Home Warranty company requirements
"""

import json
import asyncio
import streamlit as st
import pandas as pd
from datetime import datetime
from typing import Dict, List, Any
from llm_providers import get_provider, get_available_providers, OllamaProvider, AnthropicProvider

# ============================================================================
# ONTOLOGY & DATA
# ============================================================================

CMS_ONTOLOGY = {
    "capabilities": {
        "content_modeling": {
            "label": "Content Modeling",
            "facets": [
                "schema_flexibility",
                "relationship_support",
                "versioning",
                "multi_channel_support"
            ],
            "scale": "0-3",
            "importance": "critical"
        },
        "delivery": {
            "label": "Delivery & API",
            "facets": [
                "api_maturity",
                "headless_support",
                "cdn_integration",
                "performance_optimization"
            ],
            "scale": "0-3",
            "importance": "critical"
        },
        "personalization": {
            "label": "Personalization & Testing",
            "facets": [
                "native_ab_testing",
                "native_personalization",
                "segment_targeting",
                "real_time_decisioning"
            ],
            "scale": "0-3",
            "importance": "critical"
        },
        "workflow": {
            "label": "Workflow & Governance",
            "facets": [
                "content_approval_flows",
                "role_based_access",
                "versioning_control",
                "audit_logging"
            ],
            "scale": "0-3",
            "importance": "high"
        },
        "integrations": {
            "label": "Integration Ecosystem",
            "facets": [
                "crm_integration",
                "analytics_integration",
                "martech_ecosystem",
                "api_first_design"
            ],
            "scale": "0-3",
            "importance": "high"
        },
        "performance": {
            "label": "Performance & SEO",
            "facets": [
                "page_speed_optimization",
                "cdn_caching",
                "seo_metadata",
                "image_optimization"
            ],
            "scale": "0-3",
            "importance": "high"
        },
        "operational": {
            "label": "Operational & Support",
            "facets": [
                "vendor_stability",
                "documentation_quality",
                "support_tier",
                "community_maturity"
            ],
            "scale": "0-3",
            "importance": "medium"
        }
    },
    "use_cases": {
        "paid_landing_pages": {
            "label": "High-velocity paid landing pages",
            "required_capabilities": {
                "personalization": 2,
                "delivery": 2,
                "integrations": 2
            }
        },
        "enrollment_funnel": {
            "label": "Multi-step enrollment funnel with conversions",
            "required_capabilities": {
                "content_modeling": 2,
                "personalization": 3,
                "integrations": 3
            }
        },
        "multi_property": {
            "label": "Multi-property content management (3+ brand properties)",
            "required_capabilities": {
                "content_modeling": 3,
                "workflow": 2,
                "delivery": 2
            }
        },
        "consolidation": {
            "label": "Consolidate 5 legacy CMS into 3 platforms",
            "required_capabilities": {
                "integrations": 3,
                "workflow": 2,
                "operational": 2
            }
        }
    },
    "business_outcomes": {
        "conversion_lift": {
            "label": "Conversion rate lift (target: 10%+)",
            "weight": 0.25
        },
        "time_to_market": {
            "label": "Time to launch new experiments (target: <24h)",
            "weight": 0.25
        },
        "operational_efficiency": {
            "label": "Content ops efficiency (consolidation)",
            "weight": 0.20
        },
        "flexibility": {
            "label": "Flexibility for future innovation",
            "weight": 0.15
        },
        "cost_efficiency": {
            "label": "Cost efficiency & TCO",
            "weight": 0.15
        }
    }
}

PLATFORMS_DATA = {
    "HubSpot": {
        "type": "Traditional CMS + CRM",
        "category": "coupled",
        "deployment": "SaaS",
        "traffic_band": "20K+/day",
        "cost_band": "$$$",
        "summary": "Integrated CRM + CMS, strong marketing automation, limited content modeling flexibility",
        "capabilities": {
            "content_modeling": 1,
            "delivery": 1,
            "personalization": 2,
            "workflow": 2,
            "integrations": 3,
            "performance": 1,
            "operational": 3
        }
    },
    "Contentful": {
        "type": "Headless CMS",
        "category": "headless",
        "deployment": "SaaS",
        "traffic_band": "20K+/day",
        "cost_band": "$$",
        "summary": "Pure headless, rich schema, API-first, modern DX, requires front-end layer",
        "capabilities": {
            "content_modeling": 3,
            "delivery": 3,
            "personalization": 1,
            "workflow": 2,
            "integrations": 2,
            "performance": 3,
            "operational": 3
        }
    },
    "Liferay": {
        "type": "Hybrid Portal/CMS",
        "category": "hybrid",
        "deployment": "On-prem/Cloud",
        "traffic_band": "20K+/day",
        "cost_band": "$$$$",
        "summary": "Portal + CMS, enterprise workflows, steep learning curve, heavyweight",
        "capabilities": {
            "content_modeling": 2,
            "delivery": 2,
            "personalization": 2,
            "workflow": 3,
            "integrations": 2,
            "performance": 2,
            "operational": 2
        }
    },
    "Sitecore": {
        "type": "Enterprise DXP + CMS",
        "category": "monolith",
        "deployment": "On-prem/Cloud",
        "traffic_band": "20K+/day+",
        "cost_band": "$$$$$",
        "summary": "All-in-one enterprise CMS + personalization + commerce, expensive, vendor lock-in risk",
        "capabilities": {
            "content_modeling": 3,
            "delivery": 2,
            "personalization": 3,
            "workflow": 3,
            "integrations": 3,
            "performance": 2,
            "operational": 2
        }
    },
    "Sanity": {
        "type": "Headless CMS",
        "category": "headless",
        "deployment": "SaaS",
        "traffic_band": "20K+/day",
        "cost_band": "$$",
        "summary": "Rich schema + custom desk/plugins, structured content, growing ecosystem",
        "capabilities": {
            "content_modeling": 3,
            "delivery": 3,
            "personalization": 1,
            "workflow": 2,
            "integrations": 2,
            "performance": 3,
            "operational": 2
        }
    },
    "Composable (Acquia/Agility)": {
        "type": "Composable CMS",
        "category": "composable",
        "deployment": "SaaS",
        "traffic_band": "20K+/day",
        "cost_band": "$$$",
        "summary": "Best-of-breed headless + best-of-breed personalization/testing, flexible stack",
        "capabilities": {
            "content_modeling": 3,
            "delivery": 3,
            "personalization": 3,
            "workflow": 2,
            "integrations": 3,
            "performance": 3,
            "operational": 2
        }
    }
}

# ============================================================================
# STREAMLIT UI
# ============================================================================

st.set_page_config(page_title="Home Warranty CMS Evaluation", layout="wide", initial_sidebar_state="expanded")

st.title("🎯 Home Warranty CMS Evaluation Framework")
st.markdown("**Interactive tool to evaluate and score CMS platforms for consolidation strategy**")

# Sidebar: Controls
with st.sidebar:
    st.header("⚙️ Evaluation Settings")
    
    selected_use_cases = st.multiselect(
        "Select Use Cases:",
        options=list(CMS_ONTOLOGY["use_cases"].keys()),
        default=list(CMS_ONTOLOGY["use_cases"].keys()),
        format_func=lambda x: CMS_ONTOLOGY["use_cases"][x]["label"]
    )
    
    st.markdown("---")
    st.subheader("Business Driver Weights")
    
    weights = {}
    for outcome_key, outcome in CMS_ONTOLOGY["business_outcomes"].items():
        weights[outcome_key] = st.slider(
            outcome["label"],
            min_value=0.0,
            max_value=1.0,
            value=outcome["weight"],
            step=0.05,
            key=f"weight_{outcome_key}"
        )
    
    # Normalize weights
    total_weight = sum(weights.values())
    if total_weight > 0:
        weights = {k: v / total_weight for k, v in weights.items()}
    
    st.markdown("---")
    
    eval_method = st.radio(
        "Evaluation Method:",
        options=["Quick Score (Local)", "AI-Powered Analysis"],
        help="Local: Use pre-populated data. AI: Use LLM to generate detailed assessments."
    )
    
    # Provider selection (only show if AI analysis selected)
    if eval_method == "AI-Powered Analysis":
        st.markdown("---")
        st.subheader("🤖 LLM Provider")
        
        provider_type = st.radio(
            "Select Provider:",
            options=["OpenAI (Cloud)", "Ollama (Local)", "Claude (API)"],
            help="OpenAI: Cloud API (recommended). Ollama: Run locally. Claude: Anthropic API."
        )
        
        if provider_type == "Ollama (Local)":
            ollama_model = st.text_input(
                "Ollama Model:",
                value="llama3.1",
                help="Model to use (e.g., llama3.1, mistral, codellama)"
            )
            ollama_host = st.text_input(
                "Ollama Host:",
                value="http://localhost:11444",
                help="Ollama server URL"
            )
            # Check availability
            try:
                test_provider = OllamaProvider(model=ollama_model, host=ollama_host)
                if test_provider.is_available():
                    st.success(f"✅ Connected to Ollama")
                else:
                    st.warning("⚠️ Ollama running but model may need to be pulled")
            except Exception:
                st.error("❌ Cannot connect to Ollama. Is it running?")
        elif provider_type == "OpenAI (Cloud)":
            from llm_providers import OpenAIProvider
            test_provider = OpenAIProvider()
            if test_provider.is_available():
                st.success("✅ OpenAI API key configured")
            else:
                st.warning("⚠️ Set OPENAI_API_KEY in Streamlit secrets or environment")
        else:
            # Check Claude availability
            test_provider = AnthropicProvider()
            if test_provider.is_available():
                st.success("✅ Claude API key configured")
            else:
                st.warning("⚠️ Set ANTHROPIC_API_KEY environment variable")
        
        # Custom Prompt Section
        st.markdown("---")
        with st.expander("🔧 Custom Prompt (Advanced)", expanded=False):
            st.markdown("**Customize the business context for AI analysis:**")
            
            default_context = """The company is currently using HubSpot for CMS, which is essentially a CRM with an added CMS module.
HubSpot is NOT meeting expectations for improving conversions and driving enrollments - this is the PRIMARY pain point.
Currently operating across FIVE different content management systems:
- HubSpot (current primary, underperforming)
- Liferay (legacy)
- Ion (~9 years in use)
- Starmark (~9 years in use)
- Surefire (legacy)

BUSINESS GOALS:
- Primary goal: IMPROVE CONVERSION RATES
- Traffic: ~20,000 paid views/day, 6,000-7,000 unique visitors
- Need to consolidate to 3 platforms (1 won't work)

CONSTRAINTS:
- Do NOT want large enterprise CMS like Sitecore
- Do NOT want lightweight/limited/free platforms
- Considering Contentful as an option"""
            
            custom_context = st.text_area(
                "Business Context:",
                value=default_context,
                height=300,
                help="Modify this context to customize AI analysis for your specific situation"
            )
            
            st.session_state['custom_context'] = custom_context

# Main content: Tabs
tab1, tab2, tab3, tab4, tab5 = st.tabs(["📊 Platform Scores", "🌐 Live Data", "📋 Ontology", "🏗️ Recommended Stacks", "📄 Report"])

# ─────────────────────────────────────────────────────────────────────────────
# TAB 1: Platform Scores
# ─────────────────────────────────────────────────────────────────────────────

with tab1:
    st.header("Platform Capability Scores")
    
    if eval_method == "Quick Score (Local)":
        st.info("Using pre-populated capability scores. Adjust in sidebar to see impact.")
        
        # Calculate composite scores
        scores_df = []
        for platform_name, platform_data in PLATFORMS_DATA.items():
            capability_scores = platform_data["capabilities"]
            
            # Weighted score across selected use cases
            use_case_fit = 0
            for uc_key in selected_use_cases:
                uc_data = CMS_ONTOLOGY["use_cases"][uc_key]
                uc_score = 0
                count = 0
                for cap_key, required_level in uc_data["required_capabilities"].items():
                    actual_level = capability_scores.get(cap_key, 0)
                    uc_score += (actual_level / required_level) if required_level > 0 else 0
                    count += 1
                use_case_fit += (uc_score / count) if count > 0 else 0
            
            use_case_fit = (use_case_fit / len(selected_use_cases)) if selected_use_cases else 0
            
            # Business outcome fit (simplified: average capabilities)
            avg_capability = sum(capability_scores.values()) / len(capability_scores)
            business_fit = min(avg_capability / 3.0, 1.0)  # Normalize to 0-1
            
            composite_score = (use_case_fit * 0.6 + business_fit * 0.4)
            
            scores_df.append({
                "Platform": platform_name,
                "Type": platform_data["type"],
                "Category": platform_data["category"],
                "Use Case Fit": round(use_case_fit, 2),
                "Business Fit": round(business_fit, 2),
                "Composite Score": round(composite_score, 2),
                "Cost Band": platform_data["cost_band"],
                "Support": platform_data["capabilities"]["operational"]
            })
        
        scores_df = pd.DataFrame(scores_df).sort_values("Composite Score", ascending=False)
        
        # Display in columns
        col1, col2 = st.columns([2, 1])
        
        with col1:
            st.dataframe(
                scores_df.style.format({
                    "Use Case Fit": "{:.2f}",
                    "Business Fit": "{:.2f}",
                    "Composite Score": "{:.2f}"
                }).highlight_max(subset=["Composite Score"], color="lightgreen"),
                use_container_width=True
            )
        
        with col2:
            st.markdown("### Top 3 Platforms")
            for idx, row in scores_df.head(3).iterrows():
                st.metric(
                    row["Platform"],
                    f"{row['Composite Score']:.2f}",
                    f"{row['Category']}"
                )
    
    else:  # AI-Powered Analysis
        # Get provider based on sidebar selection
        if provider_type == "Ollama (Local)":
            provider = get_provider("ollama", model=ollama_model, host=ollama_host)
        elif provider_type == "OpenAI (Cloud)":
            provider = get_provider("openai")
        else:  # Claude
            provider = get_provider("anthropic")
        st.info(f"🤖 Using {provider.name} for structured assessments...")
        
        if st.button("Run AI Analysis"):
            with st.spinner(f"Analyzing with {provider.name}..."):
                try:
                    assessment_schema = {
                        "type": "object",
                        "properties": {
                            "platform": {"type": "string"},
                            "overall_fit_score": {"type": "number"},
                            "strengths": {
                                "type": "array",
                                "items": {"type": "string"}
                            },
                            "weaknesses": {
                                "type": "array",
                                "items": {"type": "string"}
                            },
                            "best_for_use_case": {"type": "string"}
                        },
                        "required": ["platform", "overall_fit_score", "strengths", "weaknesses", "best_for_use_case"]
                    }
                    
                    # Build use case context from selection
                    use_case_context = "\n".join([
                        f"- {CMS_ONTOLOGY['use_cases'][uc]['label']}"
                        for uc in selected_use_cases
                    ])
                    
                    # Get custom context from session state if available
                    business_context = st.session_state.get('custom_context', """The company is currently using HubSpot for CMS, which is essentially a CRM with an added CMS module.
HubSpot is NOT meeting expectations for improving conversions and driving enrollments - this is the PRIMARY pain point.
Currently operating across FIVE different content management systems.""")
                    
                    prompt = f"""
You are evaluating CMS platforms with this SPECIFIC business context:

## BUSINESS CONTEXT (from user input)
{business_context}

## PRIORITY USE CASES:
{use_case_context}

## PLATFORMS TO EVALUATE: {', '.join(PLATFORMS_DATA.keys())}

For each platform, provide:

1. **overall_fit_score** (0.0-1.0): Score based on:
   - Can it FIX conversion optimization problems? (weight: 30%)
   - Speed of experimentation/A/B testing (weight: 25%)  
   - Ability to consolidate legacy systems (weight: 20%)
   - Future flexibility and composability (weight: 15%)
   - Total cost of ownership (weight: 10%)

2. **strengths** (exactly 3): Specific capabilities that address the pain points described above.
   Example: "Native A/B testing allows optimizing CTAs without developer involvement"

3. **weaknesses** (exactly 3): Specific gaps relative to the stated needs.
   Example: "No built-in personalization requires integrating 3rd party CDP, adding $50K+ annual cost"

4. **best_for_use_case**: Which use case (paid_landing_pages, enrollment_funnel, multi_property_management, legacy_consolidation) is this platform BEST suited for and WHY in one sentence.

Be brutally honest. Avoid marketing language.
Return as structured JSON matching the schema provided.
"""
                    
                    response = provider.chat(prompt, assessment_schema)
                    assessments = response.content
                    
                    st.success(f"✅ AI Analysis Complete (via {response.provider})")
                    st.json(assessments)
                    
                except Exception as e:
                    st.error(f"Error calling {provider.name}: {e}")
                    if provider_type == "Ollama (Local)":
                        st.info("💡 Tip: Make sure Ollama is running (`ollama serve`) and the model is pulled (`ollama pull " + ollama_model + "`).")
                    elif provider_type == "OpenAI (Cloud)":
                        st.info("💡 Tip: Make sure your OPENAI_API_KEY is set in Streamlit secrets or as an environment variable.")
                    else:
                        st.info("💡 Tip: Make sure your ANTHROPIC_API_KEY is set as an environment variable.")

# ─────────────────────────────────────────────────────────────────────────────
# TAB 2: Live Data Fetching
# ─────────────────────────────────────────────────────────────────────────────

with tab2:
    st.header("🌐 Fetch Live Vendor Data")
    st.markdown("**Scrape real vendor documentation and extract capabilities using AI**")
    
    # Initialize session state for fetched data
    if 'live_vendor_data' not in st.session_state:
        st.session_state.live_vendor_data = {}
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        vendor_to_fetch = st.selectbox(
            "Select vendor to fetch:",
            options=["contentful", "sanity", "hubspot", "sitecore", "acquia"],
            help="Choose a vendor to scrape documentation from"
        )
    
    with col2:
        st.write("")  # Spacer
        st.write("")  # Spacer
        fetch_button = st.button("🔍 Fetch Live Data", type="primary")
    
    if fetch_button:
        try:
            from data_agents import CMSDataAgent
            
            # Progress bar with steps
            progress = st.progress(0, text="Initializing...")
            status = st.empty()
            
            # Step 1: Initialize
            status.info("🚀 Starting browser...")
            progress.progress(10, text="Launching Playwright browser...")
            
            # Get provider
            if eval_method == "AI-Powered Analysis" and provider_type == "Ollama (Local)":
                agent_provider = get_provider("ollama", model=ollama_model, host=ollama_host)
            else:
                agent_provider = get_provider("ollama")
            
            agent = CMSDataAgent(provider=agent_provider)
            
            # Step 2: Scrape
            progress.progress(25, text=f"Scraping {vendor_to_fetch} documentation...")
            status.info(f"🌐 Scraping {vendor_to_fetch} website...")
            
            raw_content = asyncio.run(agent.scrape_vendor(vendor_to_fetch))
            
            # Step 3: AI Extraction
            progress.progress(60, text="AI extracting capabilities (this takes ~30s)...")
            status.info("🤖 Ollama analyzing content...")
            
            extracted = agent.extract_capabilities(vendor_to_fetch, raw_content)
            
            # Step 4: Package results
            progress.progress(90, text="Packaging results...")
            
            from data_agents import VendorData
            data = VendorData(
                platform=vendor_to_fetch,
                capabilities=extracted.get("capabilities", {}),
                pricing_info=extracted.get("pricing_tier", "Unknown"),
                features=extracted.get("key_features", []),
                source_urls=agent.VENDOR_DOCS.get(vendor_to_fetch.lower(), []),
                raw_content=raw_content[:5000]
            )
            
            # Done
            progress.progress(100, text="✅ Complete!")
            status.success(f"✅ Fetched live data for {vendor_to_fetch}!")
            
            # Store in session state
            st.session_state.live_vendor_data[vendor_to_fetch] = data
            
            # Cleanup
            asyncio.run(agent.close())
            
        except Exception as e:
            st.error(f"Error fetching data: {e}")
            st.info("💡 Make sure Ollama is running and Playwright browsers are installed.")
    
    # Display fetched data
    if st.session_state.live_vendor_data:
        st.markdown("---")
        st.subheader("📊 Fetched Vendor Data")
        
        for vendor, data in st.session_state.live_vendor_data.items():
            with st.expander(f"**{vendor.title()}** - Live Data", expanded=True):
                col1, col2 = st.columns(2)
                
                with col1:
                    st.markdown("**Capabilities (0-3):**")
                    caps_df = pd.DataFrame([
                        {"Capability": k.replace("_", " ").title(), "Score": v}
                        for k, v in data.capabilities.items()
                    ])
                    st.dataframe(caps_df, use_container_width=True, hide_index=True)
                
                with col2:
                    st.markdown("**Key Features:**")
                    for feat in data.features[:5]:
                        st.write(f"• {feat[:100]}..." if len(feat) > 100 else f"• {feat}")
                
                st.markdown(f"**Sources:** {', '.join(data.source_urls)}")
        
        # Option to use live data in scoring
        if st.button("📥 Use Live Data in Platform Scores"):
            for vendor, data in st.session_state.live_vendor_data.items():
                vendor_key = vendor.title()
                if vendor_key in PLATFORMS_DATA:
                    PLATFORMS_DATA[vendor_key]["capabilities"] = data.capabilities
                    st.success(f"Updated {vendor_key} with live data!")
            st.info("👈 Switch to 'Platform Scores' tab to see updated rankings.")

# ─────────────────────────────────────────────────────────────────────────────
# TAB 3: Ontology Browser
# ─────────────────────────────────────────────────────────────────────────────

with tab3:
    st.header("CMS Ontology")
    st.markdown("**Structured capability model used for evaluation**")
    
    ontology_view = st.radio("View:", ["Capabilities", "Use Cases", "Business Outcomes"])
    
    if ontology_view == "Capabilities":
        for cap_key, cap_data in CMS_ONTOLOGY["capabilities"].items():
            with st.expander(f"🔹 {cap_data['label']} (Importance: {cap_data['importance']})"):
                st.markdown(f"**Scale:** {cap_data['scale']}")
                st.markdown(f"**Facets:**")
                for facet in cap_data["facets"]:
                    st.write(f"  • {facet.replace('_', ' ').title()}")
    
    elif ontology_view == "Use Cases":
        for uc_key, uc_data in CMS_ONTOLOGY["use_cases"].items():
            with st.expander(f"📌 {uc_data['label']}"):
                st.markdown("**Required Capabilities:**")
                for cap, level in uc_data["required_capabilities"].items():
                    cap_name = CMS_ONTOLOGY["capabilities"][cap]["label"]
                    st.write(f"  • {cap_name}: Level {level}/3")
    
    else:  # Business Outcomes
        outcome_weights_df = pd.DataFrame([
            {
                "Outcome": outcome["label"],
                "Weight": outcome["weight"],
                "Normalized": f"{(outcome['weight'] / sum(o['weight'] for o in CMS_ONTOLOGY['business_outcomes'].values())) * 100:.1f}%"
            }
            for outcome in CMS_ONTOLOGY["business_outcomes"].values()
        ])
        st.dataframe(outcome_weights_df, use_container_width=True)

# ─────────────────────────────────────────────────────────────────────────────
# TAB 4: AI-Driven Recommended Stacks
# ─────────────────────────────────────────────────────────────────────────────

with tab4:
    st.header("🏗️ AI-Recommended Architecture Stacks")
    st.markdown("**Get personalized stack recommendations based on your specific requirements**")
    
    # Initialize session state for stack recommendations
    if 'stack_recommendations' not in st.session_state:
        st.session_state.stack_recommendations = None
    
    # Get provider type - default to OpenAI if not in AI mode
    current_provider = provider_type if 'provider_type' in dir() else "OpenAI (Cloud)"
    st.info(f"🤖 Using {current_provider} for stack recommendations")
    
    if st.button("🚀 Generate AI Stack Recommendations", type="primary"):
        with st.spinner("Analyzing optimal technology stacks..."):
            try:
                # Get provider based on sidebar selection
                if 'provider_type' in dir() and provider_type == "Ollama (Local)":
                    stack_provider = get_provider("ollama", model=ollama_model, host=ollama_host)
                elif 'provider_type' in dir() and provider_type == "OpenAI (Cloud)":
                    stack_provider = get_provider("openai")
                elif 'provider_type' in dir():
                    stack_provider = get_provider("anthropic")
                else:
                    stack_provider = get_provider("openai")  # Default
                
                stack_schema = {
                    "type": "object",
                    "properties": {
                        "recommended_stacks": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "name": {"type": "string"},
                                    "fit_score": {"type": "number"},
                                    "components": {"type": "array", "items": {"type": "string"}},
                                    "pros": {"type": "array", "items": {"type": "string"}},
                                    "cons": {"type": "array", "items": {"type": "string"}},
                                    "migration_approach": {"type": "string"},
                                    "timeline_months": {"type": "integer"},
                                    "cost_tier": {"type": "string"},
                                    "best_for": {"type": "string"}
                                }
                            }
                        },
                        "top_recommendation": {"type": "string"},
                        "migration_strategy": {"type": "string"}
                    }
                }
                
                # Build use case context
                use_case_labels = [CMS_ONTOLOGY['use_cases'][uc]['label'] for uc in selected_use_cases]
                
                stack_prompt = f"""
Based on CINCH's specific situation, recommend 3 technology stack options:

## CINCH CONTEXT
- Currently on HubSpot CMS which is FAILING at conversion optimization
- Operating across 5 legacy CMS: HubSpot, Liferay, Ion (~9yr), Starmark (~9yr), Surefire
- Need to consolidate to 3 platforms (1 won't work)
- Traffic: 20K paid views/day, 6-7K unique visitors
- Primary goal: IMPROVE CONVERSION RATES
- Avoid: Sitecore (too expensive), Liferay (too lightweight)
- Considering: Contentful as an option

## PRIORITY USE CASES
{chr(10).join(['- ' + uc for uc in use_case_labels])}

## AVAILABLE PLATFORMS TO INCLUDE
{', '.join(PLATFORMS_DATA.keys())}

For each of the 3 stack options, provide:
1. **name**: Descriptive stack name (e.g., "Headless + CDP Stack")
2. **fit_score**: 0.0-1.0 based on CINCH's requirements
3. **components**: List of 3-5 specific technologies (e.g., ["Contentful", "Segment CDP", "Next.js", "HubSpot CRM"])
4. **pros**: 3 specific benefits addressing HubSpot's failures
5. **cons**: 2-3 honest drawbacks
6. **migration_approach**: "strangler_fig", "phased", or "big_bang"
7. **timeline_months**: Realistic implementation timeline
8. **cost_tier**: "$", "$$", or "$$$"
9. **best_for**: Which of CINCH's use cases this stack excels at

Also provide:
- **top_recommendation**: Which stack you recommend and why (1-2 sentences)
- **migration_strategy**: Recommended approach to move from 5 CMS to 3 platforms (2-3 sentences)

Be specific to CINCH. Consider their HubSpot pain points.
Return as structured JSON.
"""
                
                response = stack_provider.chat(stack_prompt, stack_schema)
                st.session_state.stack_recommendations = response.content
                st.success("✅ Stack recommendations generated!")
                
            except Exception as e:
                st.error(f"Error generating recommendations: {e}")
                if provider_type == "OpenAI (Cloud)":
                    st.info("💡 Tip: Make sure your OPENAI_API_KEY is set in Streamlit secrets.")
    
    # Display recommendations if available
    if st.session_state.stack_recommendations:
        recs = st.session_state.stack_recommendations
        
        # Top recommendation banner
        if "top_recommendation" in recs:
            st.success(f"🏆 **Top Recommendation:** {recs['top_recommendation']}")
        
        if "migration_strategy" in recs:
            st.info(f"📋 **Migration Strategy:** {recs['migration_strategy']}")
        
        st.markdown("---")
        
        # Display each stack
        for i, stack in enumerate(recs.get("recommended_stacks", [])[:3]):
            st.markdown(f"### Option {chr(65+i)}: {stack.get('name', 'Stack ' + str(i+1))}")
            
            col1, col2, col3 = st.columns([2, 1, 1])
            
            with col1:
                st.markdown("**Components:**")
                for comp in stack.get("components", []):
                    st.write(f"  • {comp}")
                
                st.markdown("**Pros:**")
                for pro in stack.get("pros", []):
                    st.write(f"  ✅ {pro}")
                
                st.markdown("**Cons:**")
                for con in stack.get("cons", []):
                    st.write(f"  ❌ {con}")
            
            with col2:
                st.metric("Fit Score", f"{stack.get('fit_score', 0):.2f}")
                st.metric("Timeline", f"{stack.get('timeline_months', '?')} months")
                st.metric("Cost Tier", stack.get("cost_tier", "$$"))
            
            with col3:
                st.markdown("**Migration:**")
                approach = stack.get("migration_approach", "phased")
                if approach == "strangler_fig":
                    st.write("🌿 Strangler Fig (gradual)")
                elif approach == "big_bang":
                    st.write("💥 Big Bang (all at once)")
                else:
                    st.write("📦 Phased (batch cutover)")
                
                st.markdown("**Best For:**")
                st.write(stack.get("best_for", "General use"))
            
            st.markdown("---")
    else:
        st.info("👆 Click the button above to generate AI-powered stack recommendations based on your specific requirements.")

# ─────────────────────────────────────────────────────────────────────────────
# TAB 5: Report & Export
# ─────────────────────────────────────────────────────────────────────────────

with tab5:
    st.header("📄 Evaluation Report & Export")
    
    # Generate report content
    report_content = f"""
# Home Warranty CMS Evaluation Report
**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

## Executive Summary
The company is evaluating headless and composable CMS solutions to consolidate 5 legacy platforms (HubSpot, Liferay, Ion, Starmark, Surefire) into 3 unified platforms. The primary goal is to improve conversion optimization and enrollment rates for a site receiving ~20K paid views/day.

## Current State
- **Primary CMS:** HubSpot (underperforming for conversions)
- **Legacy Systems:** Liferay, Ion (~9 years), Starmark (~9 years), Surefire
- **Traffic:** ~20,000 paid views/day, 6,000-7,000 unique visitors
- **Goal:** Consolidate to 3 platforms, improve conversion rates

## Selected Use Cases
{chr(10).join(['- ' + CMS_ONTOLOGY['use_cases'][uc]['label'] for uc in selected_use_cases])}

## Key Findings
- **Best overall fit:** Composable CMS + HubSpot CRM
- **Best headless-only fit:** Contentful
- **Best for quick wins:** HubSpot + Headless hybrid

## Recommendation
Pursue a **phased migration** approach using the Strangler Fig pattern to gradually replace legacy systems while maintaining business continuity.
"""
    
    st.markdown(report_content)
    
    st.markdown("---")
    st.subheader("📥 Export Options")
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        if st.button("📄 Export to DOCX", type="primary"):
            try:
                from docx import Document
                from docx.shared import Inches, Pt
                import io
                
                doc = Document()
                doc.add_heading('Home Warranty CMS Evaluation Report', 0)
                doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
                
                doc.add_heading('Executive Summary', level=1)
                doc.add_paragraph('The company is evaluating headless and composable CMS solutions to consolidate 5 legacy platforms into 3 unified platforms.')
                
                doc.add_heading('Current State', level=1)
                doc.add_paragraph('• Primary CMS: HubSpot (underperforming)')
                doc.add_paragraph('• Legacy Systems: Liferay, Ion, Starmark, Surefire')
                doc.add_paragraph('• Traffic: ~20K paid views/day')
                
                doc.add_heading('Selected Use Cases', level=1)
                for uc in selected_use_cases:
                    doc.add_paragraph(f"• {CMS_ONTOLOGY['use_cases'][uc]['label']}", style='List Bullet')
                
                doc.add_heading('Recommendation', level=1)
                doc.add_paragraph('Pursue a phased migration using the Strangler Fig pattern.')
                
                buffer = io.BytesIO()
                doc.save(buffer)
                buffer.seek(0)
                
                st.download_button(
                    label="⬇️ Download DOCX",
                    data=buffer,
                    file_name="cms_evaluation_report.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )
            except ImportError:
                st.error("python-docx not installed")
    
    with col2:
        if st.button("📊 Export to PPTX", type="primary"):
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                import io
                
                prs = Presentation()
                
                # Title slide
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = "Home Warranty CMS Evaluation"
                slide.placeholders[1].text = f"Generated: {datetime.now().strftime('%Y-%m-%d')}"
                
                # Executive Summary slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Executive Summary"
                slide.placeholders[1].text = "• Consolidating 5 CMS platforms to 3\n• Primary goal: Improve conversion rates\n• Traffic: 20K paid views/day\n• Current CMS (HubSpot) underperforming"
                
                # Key Findings slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Key Findings"
                slide.placeholders[1].text = "• Best overall: Composable CMS + HubSpot CRM\n• Best headless: Contentful\n• Best quick wins: HubSpot + Headless hybrid"
                
                # Recommendation slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Recommendation"
                slide.placeholders[1].text = "• Use Strangler Fig migration pattern\n• Phased approach over 12-18 months\n• Start with highest-traffic landing pages"
                
                buffer = io.BytesIO()
                prs.save(buffer)
                buffer.seek(0)
                
                st.download_button(
                    label="⬇️ Download PPTX",
                    data=buffer,
                    file_name="cms_evaluation_report.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except ImportError:
                st.error("python-pptx not installed")
    
    with col3:
        if st.button("📑 Export to PDF", type="primary"):
            try:
                from reportlab.lib.pagesizes import letter
                from reportlab.pdfgen import canvas
                from reportlab.lib.units import inch
                import io
                
                buffer = io.BytesIO()
                c = canvas.Canvas(buffer, pagesize=letter)
                width, height = letter
                
                # Title
                c.setFont("Helvetica-Bold", 20)
                c.drawString(1*inch, height - 1*inch, "Home Warranty CMS Evaluation Report")
                
                c.setFont("Helvetica", 12)
                c.drawString(1*inch, height - 1.5*inch, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
                
                # Executive Summary
                c.setFont("Helvetica-Bold", 14)
                c.drawString(1*inch, height - 2.5*inch, "Executive Summary")
                c.setFont("Helvetica", 11)
                c.drawString(1*inch, height - 3*inch, "Consolidating 5 CMS platforms to 3 unified platforms.")
                c.drawString(1*inch, height - 3.3*inch, "Primary goal: Improve conversion rates by 10%+")
                c.drawString(1*inch, height - 3.6*inch, "Traffic: ~20,000 paid views/day")
                
                # Key Findings
                c.setFont("Helvetica-Bold", 14)
                c.drawString(1*inch, height - 4.5*inch, "Key Findings")
                c.setFont("Helvetica", 11)
                c.drawString(1*inch, height - 5*inch, "• Best overall: Composable CMS + HubSpot CRM")
                c.drawString(1*inch, height - 5.3*inch, "• Best headless: Contentful")
                c.drawString(1*inch, height - 5.6*inch, "• Best quick wins: HubSpot + Headless hybrid")
                
                # Recommendation
                c.setFont("Helvetica-Bold", 14)
                c.drawString(1*inch, height - 6.5*inch, "Recommendation")
                c.setFont("Helvetica", 11)
                c.drawString(1*inch, height - 7*inch, "Pursue phased migration using Strangler Fig pattern")
                
                c.save()
                buffer.seek(0)
                
                st.download_button(
                    label="⬇️ Download PDF",
                    data=buffer,
                    file_name="cms_evaluation_report.pdf",
                    mime="application/pdf"
                )
            except ImportError:
                st.error("reportlab not installed")

# ─────────────────────────────────────────────────────────────────────────────
# Footer
# ─────────────────────────────────────────────────────────────────────────────

st.markdown("---")
st.markdown("""
**Built with:** Python, Streamlit, OpenAI GPT-4  
**Ontology Version:** 1.0  
**Last Updated:** 2025-12-11
""")
